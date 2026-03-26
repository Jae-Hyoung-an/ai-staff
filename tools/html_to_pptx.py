"""
신규 지역 진출 기획 SOP — HTML 슬라이드를 편집 가능한 PPTX로 변환
python-pptx로 네이티브 요소(텍스트박스, 도형, 테이블) 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0B, 0x1F, 0x3B)
GREEN = RGBColor(0x1B, 0x73, 0x40)
ORANGE = RGBColor(0xC7, 0x5B, 0x12)
GRAY = RGBColor(0x6B, 0x7A, 0x90)
TEXT_DARK = RGBColor(0x34, 0x48, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
BORDER_COLOR = RGBColor(0xE3, 0xEA, 0xF4)
CARD_BORDER = RGBColor(0xD7, 0xE2, 0xF0)

SLIDE_W = Emu(12192000)  # 1280px at 96dpi ≈ 13.33in
SLIDE_H = Emu(6858000)   # 720px at 96dpi ≈ 7.5in

FONT_NAME = 'Noto Sans KR'


def px(val):
    return Emu(int(val * 9525))


def add_top_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(0), px(0), px(1280), px(18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    return shape


def add_left_rail(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(0), px(18), px(8), px(702))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=NAVY, alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card_bg(slide, left, top, width, height, top_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(left), px(top), px(width), px(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1.5)
    if top_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(left), px(top), px(width), px(7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = top_color
        bar.line.fill.background()
    return shape


def add_multi_text(slide, left, top, width, height, lines, default_size=14):
    """lines: list of (text, font_size, bold, color)"""
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(2)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=TEXT_DARK):
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.level = 0
        p.space_after = Pt(2)
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        existing = pPr.findall(qn('a:buChar'))
        for e in existing:
            pPr.remove(e)
        pPr.append(buChar)
    return txBox


def add_purpose_card(slide, top_y, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(top_y), px(1120), px(90))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(1)

    add_multi_text(slide, 150, top_y + 16, 1048, 62, [
        ("목적", 14, True, NAVY),
        (text, 18, True, NAVY),
    ])


def add_org_card(slide, left, top, width, height, top_color, title, subtitle, bullets):
    add_card_bg(slide, left, top, width, height, top_color)
    add_multi_text(slide, left + 24, top + 20, width - 48, 50, [
        (title, 18, True, NAVY),
        (subtitle, 12, True, GRAY),
    ])
    add_bullet_list(slide, left + 24, top + 76, width - 48, height - 90, bullets)


def add_bottom_bar(slide, top_y, left_label, left_text, right_label, right_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(top_y), px(1120), px(94))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(1)

    add_multi_text(slide, 120, top_y + 14, 520, 60, [
        (left_label, 13, True, NAVY),
        (left_text, 13, False, TEXT_DARK),
    ])
    add_multi_text(slide, 656, top_y + 14, 544, 60, [
        (right_label, 13, True, NAVY),
        (right_text, 13, False, TEXT_DARK),
    ])


def add_footer(slide, text, top_y=694):
    add_textbox(slide, 96, top_y, 1120, 20, text, font_size=11, color=GRAY)


# ─── Slide builders ───

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_top_bar(slide)
    add_left_rail(slide)

    # decorative bars
    for (l, t, w, op) in [(1040, 86, 180, 1), (1090, 112, 130, 0.6), (1140, 138, 80, 0.4)]:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(l), px(t), px(w), px(10))
        s.fill.solid()
        s.fill.fore_color.rgb = NAVY
        s.line.fill.background()

    add_textbox(slide, 96, 210, 1088, 70, "신규 지역 진출 기획 SOP",
                font_size=48, bold=True, color=NAVY)
    add_textbox(slide, 96, 290, 1088, 50,
                "권역/세트분배 기능 기반으로 지역 확정부터 영업·공급 목표 수립까지의 기획 프로세스를 표준화",
                font_size=18, color=TEXT_DARK)

    # divider
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(96), px(392), px(640), px(2))
    s.fill.solid(); s.fill.fore_color.rgb = NAVY; s.line.fill.background()

    add_multi_text(slide, 96, 418, 1088, 60, [
        ("작성일", 15, True, NAVY),
        ("2026.03.24", 15, False, TEXT_DARK),
    ])

    add_footer(slide, "적용 기간: 런칭 L-4주 ~ L-1주 (런칭 이후는 B존 TF 및 주간 S&OP로 전환)", 650)


def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide); add_left_rail(slide)

    add_textbox(slide, 96, 44, 1120, 60, "개요: 목적 · 적용 범위 · 참여 조직",
                font_size=30, bold=True, color=NAVY)

    # Purpose section
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(120), px(1120), px(148))
    s.fill.solid(); s.fill.fore_color.rgb = LIGHT_BG
    s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)

    add_textbox(slide, 120, 140, 260, 30, "1) 목적", font_size=15, bold=True, color=NAVY)
    add_multi_text(slide, 120, 172, 1072, 80, [
        ("신규 지역 진출의 표준 기획 프로세스 정립", 18, True, NAVY),
        ("지역 확정 → 권역(Zone) 설계 → 원/판가 설정 → 영업 목표물량·기사 공급목표 수립까지", 13, False, TEXT_DARK),
    ])

    # Scope section
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(286), px(1120), px(132))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)

    add_textbox(slide, 120, 306, 260, 30, "2) 적용 범위", font_size=15, bold=True, color=NAVY)
    add_multi_text(slide, 150, 344, 1042, 64, [
        ("런칭 L-4주 ~ L-1주 (런칭 전)", 14, True, NAVY),
        ("런칭 이후에는 B존 TF 워크플로우로 전환, 초기 기획 이행 점검·조정은 주간 S&OP 미팅(별도 SOP)으로 운영", 13, False, TEXT_DARK),
    ])

    # Orgs
    add_textbox(slide, 96, 444, 1120, 30, "3) 참여 조직 및 역할 (R&R)", font_size=15, bold=True, color=NAVY)

    add_org_card(slide, 96, 490, 360, 190, NAVY,
                 "기획조직", "조율 및 플래닝 하달",
                 ["전체 기획 사이클 총괄 및 일정/산출물 관리",
                  "데이터 제공, 단계별 리뷰·승인",
                  "진출 목표(물량·공급) 최종 확정"])

    add_org_card(slide, 476, 490, 360, 190, GREEN,
                 "영업조직", "영업 담당",
                 ["권역(Zone) 기획 및 시장/경쟁사 조사",
                  "판가 설정 및 영업 목표물량 수립",
                  "상점 파이프라인 관리"])

    add_org_card(slide, 856, 490, 360, 190, ORANGE,
                 "운영조직", "공급 담당",
                 ["협력사(벤더) 확보 가능성 검토",
                  "기사 원가·보상 설계 및 공급 기획",
                  "배송 Capa 산출 (세트/프렌즈)"])

    add_footer(slide, "조직별 역할은 단계별로 '주도/협의/검토'가 달라지며, 본 SOP는 런칭 전 4주 기획 구간의 표준 흐름을 정의합니다.")


def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide); add_left_rail(slide)

    add_textbox(slide, 96, 48, 1120, 60, "전체 프로세스 흐름 (L-4주 → 런칭)",
                font_size=30, bold=True, color=NAVY)

    # Legend bar
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(106), px(1120), px(36))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)

    for (label, color_rgb, x) in [("기획(주도)", NAVY, 132), ("영업(주도)", GREEN, 268), ("운영(주도/검토)", ORANGE, 404)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x - 18), px(118), px(12), px(12))
        dot.fill.solid(); dot.fill.fore_color.rgb = color_rgb; dot.line.fill.background()
        add_textbox(slide, x, px(110)//9525, 140, 24, label, font_size=12, bold=True, color=TEXT_DARK)

    add_textbox(slide, 820, 110, 380, 24, "Step 4-1/4-2는 병렬 진행",
                font_size=12, bold=True, color=GRAY, alignment=PP_ALIGN.RIGHT)

    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(120), px(208), px(1040), px(2))
    line.fill.solid(); line.fill.fore_color.rgb = CARD_BORDER; line.line.fill.background()

    # Step cards on timeline
    steps = [
        (96, "Step 1. 진출 지역 확정", "L-4주 · 기획 주도", "산출물: 진출 대상 지역 확정안", NAVY),
        (336, "Step 2. 권역 기획", "L-4~3주 · 영업 주도", "산출물: 권역 맵 · Zone 목록 · 관리정책", GREEN),
        (576, "Step 3. 원/판가·기준", "L-3~2.5주 · 영업+운영", "산출물: 원/판가표 · 영업기준표 · 개시일자", NAVY),
    ]

    for (left, title, sub, output, color) in steps:
        # node
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(left + 54), px(198), px(22), px(22))
        node.fill.solid(); node.fill.fore_color.rgb = color; node.line.fill.background()
        # card
        add_card_bg(slide, left, px(240)//9525, 230, 150, color)
        add_multi_text(slide, left + 16, px(258)//9525, 198, 42, [
            (title, 14, True, NAVY),
            (sub, 11, True, GRAY),
        ])
        add_multi_text(slide, left + 16, px(310)//9525, 198, 60, [
            (output, 12, False, TEXT_DARK),
        ])

    # Step 4 (larger)
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(870), px(198), px(22), px(22))
    node.fill.solid(); node.fill.fore_color.rgb = NAVY; node.line.fill.background()
    add_card_bg(slide, 816, px(240)//9525, 240, 270, NAVY)
    add_multi_text(slide, 834, px(258)//9525, 206, 40, [
        ("Step 4. 목표 확정 (병렬)", 14, True, NAVY),
        ("L-2.5~1주 · 기획 주도", 11, True, GRAY),
    ])

    # 4-1 / 4-2 inner cards
    for (y, title, desc) in [(304, "4-1 영업 목표물량", "Zone×주차 목표 오더"),
                              (384, "4-2 기사 공급목표", "세트 수 + 프렌즈 목표")]:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(834), px(y), px(206), px(68))
        s.fill.solid(); s.fill.fore_color.rgb = LIGHT_BG
        s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)
        add_multi_text(slide, 856, y + 12, 170, 44, [
            (title, 13, True, NAVY),
            (desc, 11, False, TEXT_DARK),
        ])

    add_multi_text(slide, 834, 460, 206, 40, [
        ("산출물", 12, True, NAVY),
        ("Zone×주차 목표(오더/세트/프렌즈)", 11, False, TEXT_DARK),
    ])

    # Launch node
    add_card_bg(slide, 1066, px(240)//9525, 150, 150, GRAY)
    add_multi_text(slide, 1082, px(258)//9525, 120, 40, [
        ("런칭", 14, True, NAVY),
        ("서비스 개시", 11, True, GRAY),
    ])
    add_multi_text(slide, 1082, px(310)//9525, 120, 60, [
        ("운영 전환", 12, True, NAVY),
        ("B존 TF / 주간 S&OP(별도 SOP)", 11, False, TEXT_DARK),
    ])

    # System linkage section
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(540), px(1120), px(130))
    s.fill.solid(); s.fill.fore_color.rgb = LIGHT_BG
    s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)

    add_textbox(slide, 120, 555, 1080, 24, "시스템 연동 포인트 (런칭 전)", font_size=14, bold=True, color=NAVY)

    sys_items = [
        (120, "Step 2: 인트라 권역관리", "관리정책 생성 + 연결 Zone 설정"),
        (488, "Step 3: 인트라 권역관리", "기사 원가 요금제 + 보상기준 입력(전체 저장)"),
        (856, "Step 4-2: 세트분배 관리", "Zone-벤더 연결 + 벤더별 세트 수 입력"),
    ]
    for (x, title, desc) in sys_items:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(585), px(352), px(68))
        s.fill.solid(); s.fill.fore_color.rgb = WHITE
        s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)
        add_multi_text(slide, x + 16, 597, 320, 44, [
            (title, 12, True, NAVY),
            (desc, 11, False, TEXT_DARK),
        ])

    add_footer(slide, "본 슬라이드는 런칭 전 4주 기획 구간의 표준 흐름(산출물·시스템 연동)을 한 장으로 요약합니다.", 688)


def build_step_slide(prs, title, subtitle, purpose_text,
                     orgs, bottom_left, bottom_right, footer_text,
                     org_colors=None):
    """Generic step detail slide (pages 4-8)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide); add_left_rail(slide)

    add_textbox(slide, 96, 46, 1120, 50, title, font_size=30, bold=True, color=NAVY)
    add_textbox(slide, 96, 92, 1120, 26, subtitle, font_size=13, bold=True, color=GRAY)
    add_purpose_card(slide, 132, purpose_text)

    add_textbox(slide, 96, 242, 1120, 24, "조직별 역할 (R&R)", font_size=15, bold=True, color=NAVY)

    if org_colors is None:
        org_colors = [NAVY, GREEN, ORANGE]

    for i, (org_title, org_sub, org_bullets) in enumerate(orgs):
        left = 96 + i * 380
        add_org_card(slide, left, 276, 360, 300, org_colors[i],
                     org_title, org_sub, org_bullets)

    add_bottom_bar(slide, 598, bottom_left[0], bottom_left[1],
                   bottom_right[0], bottom_right[1])
    add_footer(slide, footer_text)


def build_slide4(prs):
    build_step_slide(prs,
        "Step 1. 진출 지역 확정",
        "L-4주 킥오프 | 기획 주도",
        "사전 후보 중 진출 대상 지역을 최종 확정하고 기획을 착수",
        [
            ("기획 (주도)", "Kick-off & 의사결정",
             ["킥오프 미팅 주재", "진출 지역 최종 확정", "의사결정 및 공지"]),
            ("영업 (의견)", "시장/파이프라인 관점",
             ["상점 영업 파이프라인 상황 공유", "시장 관점의 리스크/기회 의견 제시"]),
            ("운영 (의견)", "공급 실현성 관점",
             ["벤더/프렌즈 수급 가능성 의견", "동선·실행 가능성 관점의 리스크 제시"]),
        ],
        ("산출물", "진출 대상 지역 확정안"),
        ("완료 기준", "킥오프에서 지역 확정 및 다음 단계 착수"),
        "PoC 단계에서는 후보지역이 사전 선정되어 있으며, 킥오프에서 최종 확정합니다.",
        [NAVY, GREEN, ORANGE])


def build_slide5(prs):
    build_step_slide(prs,
        "Step 2. 권역 기획",
        "L-4~3주 | 영업 주도",
        "목표지역을 상권/거주지 특성에 맞는 권역(Zone)으로 분할",
        [
            ("영업 (주도)", "Zone 설계 (Sales First)",
             ["경쟁사 권역 구조/커버리지 분석",
              "Zone 초안 설계 (상점 밀집도·상권 특성 반영)",
              "인트라 등록용 관리정책 구조 설계"]),
            ("운영 (검토)", "공급/동선 실현성",
             ["협력사(벤더) 확보 가능성 검토",
              "배송 Capa·동선 효율성 검증",
              "공급 관점 Zone 조정 피드백"]),
            ("기획 (조율)", "기준 제시 & 최종 승인",
             ["Zone 설계 기준 가이드 제공",
              "영업/운영 의견 종합",
              "Zone 최종 확정 승인"]),
        ],
        ("산출물", "권역 구획 맵 · Zone 목록표 · 관리정책(Zone 연결 완료)"),
        ("시스템", "인트라 권역관리 → 관리정책 생성 + 연결 Zone 등록"),
        "관리정책은 Step 2에서 생성/Zone 연결까지 진행하며, 보상·원가·세트 물량 입력은 후속 단계에서 설정합니다.",
        [GREEN, ORANGE, NAVY])


def build_slide6(prs):
    build_step_slide(prs,
        "Step 3. 원/판가 + 영업기준 기획",
        "L-3~2.5주 | 영업(판가) + 운영(원가) + 영업기준 협의",
        "권역별 영업이익 목표를 고려하여 상점 판가와 기사 원가 설정",
        [
            ("영업 (판가 주도)", "시장·가격 전략",
             ["경쟁사 배달비 조사", "상점 지불의사 조사",
              "영업이익 목표 반영", "권역별 판가 확정"]),
            ("운영 (원가 주도)", "요금·보상 설계",
             ["기사 원가 산출(판가에서 역산)",
              "원가 요금제 설계(50m당·픽업·배송완료)",
              "보상기준 설계(수락률·성공 슬롯)", "수익성 검증"]),
            ("기획 (조율)", "정합성·손익 검증",
             ["원/판가 정합성 리뷰(판가-원가=영업이익)",
              "P&L 시뮬레이션(시나리오별 손익)",
              "승인 및 의사결정 지원"]),
        ],
        ("산출물", "원/판가표, 수익성 시뮬레이션, 인트라 원가/보상, Zone별 영업 기준표, 영업 개시 일자"),
        ("시스템", "인트라 권역관리 원가/보상 입력 · 영업기준표 · 영업 개시 일자 확정"),
        "영업 기준표(Zone별 운영시간·배송 권역)와 영업 개시 일자는 기획+영업+운영 협의 후 기획이 확정합니다.",
        [GREEN, ORANGE, NAVY])


def build_slide7(prs):
    build_step_slide(prs,
        "Step 4-1. 영업 목표물량",
        "L-2.5~1주 | 기획 주도 | ※ Step 4-2와 병렬 진행",
        "사업계획 기반으로 Zone×주차별 영업 목표(오더 건수) 확정",
        [
            ("기획 (주도)", "목표 설정 & 확정",
             ["사업계획 기반 목표 오더 초안 작성",
              "Zone별·주차별 목표 하달",
              "영업 의견 반영 후 최종 확정"]),
            ("영업 (협의)", "현장/파이프라인 반영",
             ["상점 파이프라인 현황 기반 조정",
              "시장 상황/경쟁 구도 반영 의견",
              "Zone·주차별 목표 조정안 제시"]),
            ("운영 (참고)", "Capa 관점 의견",
             ["목표 물량 대비 배송 Capa 의견",
              "공급 계획과의 정합성 확인",
              "Step 4-2(공급목표)와 연동 피드백"]),
        ],
        ("산출물", "Zone×주차별 영업 목표(오더)"),
        ("프로세스", "기획 초안 하달 → 영업 조정 의견 → 협의 → 기획 최종 확정"),
        "Step 4-1과 Step 4-2는 동일 기간에 병렬 진행되며, 목표 물량과 공급 계획의 정합성을 함께 점검합니다.",
        [NAVY, GREEN, ORANGE])


def build_slide8(prs):
    build_step_slide(prs,
        "Step 4-2. 기사 공급목표",
        "L-2.5~1주 | 기획 주도 | ※ Step 4-1과 병렬 진행",
        "Zone×주차별 벤더 세트 수 및 활성 프렌즈 인원 목표 확정",
        [
            ("기획 (주도)", "공급목표 초안 · 확정 · 완료 선언",
             ["데이터 참고자료(물량 트렌드·프렌즈 생산성) 기반",
              "Zone별·주차별 세트 수 + 프렌즈 수 초안 하달",
              "운영 의견 반영 후 최종 확정",
              "진출 기획 완료 선언"]),
            ("운영 (협의)", "공급 실현성 검토 · 실행",
             ["벤더 확보 가능성 검증", "프렌즈 수급 가능성 검증",
              "조정 의견 제시", "벤더 계약 조율",
              "프렌즈 확보 계획 수립"]),
            ("영업 (참고)", "물량 리스크 관점 (4-1 연동)",
             ["영업 파이프라인 대비 공급 계획 적절성 의견",
              "목표 물량 리스크 조기 공유 (Step 4-1 연동)"]),
        ],
        ("산출물", "세트 수/프렌즈 목표, 벤더 계약 계획, 프렌즈 확보 계획, 인트라 세트분배 설정"),
        ("시스템", "세트분배 관리 → Zone-벤더 연결 + 벤더별 세트 수 입력 + 배차 엔진 연동"),
        "Step 4-1(영업 목표물량)과 Step 4-2(기사 공급목표)는 병렬로 진행하며, 모두 확정 시 '진출 기획 완료'를 선언합니다.",
        [NAVY, ORANGE, GREEN])


def build_slide9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide); add_left_rail(slide)

    add_textbox(slide, 96, 38, 1120, 52, "단계별 산출물 요약 & 런칭 전환",
                font_size=30, bold=True, color=NAVY)

    # Table
    rows = 7  # header + 6 data rows
    cols = 4
    tbl_shape = slide.shapes.add_table(rows, cols, px(96), px(100), px(1120), px(390))
    table = tbl_shape.table

    table.columns[0].width = px(208)
    table.columns[1].width = px(564)
    table.columns[2].width = px(160)
    table.columns[3].width = px(188)

    headers = ["단계", "최종 산출물", "확정 주체", "시점"]
    data = [
        ["Step 1\n진출 지역 확정", "진출 대상 지역 확정안", "기획", "L-4주 킥오프"],
        ["Step 2\n권역 기획", "권역 맵, Zone 목록, 관리정책(Zone 연결)", "기획 승인", "L-3주"],
        ["Step 3\n원/판가+영업기준", "원/판가표, 수익성 시뮬레이션, 원가/보상, 영업 기준표, 영업 개시 일자", "기획 승인", "L-2.5주"],
        ["Step 4-1\n영업 목표물량", "Zone×주차별 영업 목표(오더)", "기획 확정", "L-1주 이전"],
        ["Step 4-2\n기사 공급목표", "세트 수/프렌즈 목표, 벤더 계약, 프렌즈 확보 계획, 세트분배 설정", "기획 확정", "L-1주 이전"],
        ["진출 기획 완료", "Step 1~4 통합 최종 산출물 패키지 (런칭 준비본)", "기획 최종 확정", "L-1주 완료 미팅"],
    ]

    def set_cell(cell, text, bold=False, size=12, color=TEXT_DARK):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        set_cell(cell, h, bold=True, size=13, color=NAVY)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BG

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            is_bold = (c == 0)
            color = NAVY if c == 0 else TEXT_DARK
            set_cell(cell, val, bold=is_bold, size=12, color=color)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFE)

    # Launch transition
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(96), px(506), px(1120), px(108))
    s.fill.solid(); s.fill.fore_color.rgb = LIGHT_BG
    s.line.color.rgb = BORDER_COLOR; s.line.width = Pt(1)

    add_textbox(slide, 164, 522, 1028, 28, "런칭 및 운영 전환", font_size=15, bold=True, color=NAVY)
    add_textbox(slide, 164, 556, 1028, 48,
                "서비스 런칭 → B존 TF 워크플로우로 전환 · 초기 기획 이행 점검/조정은 주간 S&OP 미팅(별도 SOP)으로 운영",
                font_size=13, color=TEXT_DARK)

    add_footer(slide, "표는 런칭 전 4주 기획 구간에서 '무엇을 언제 확정하는지'를 산출물 기준으로 정리합니다.", 626)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)
    build_slide5(prs)
    build_slide6(prs)
    build_slide7(prs)
    build_slide8(prs)
    build_slide9(prs)

    out_dir = os.path.join(os.path.dirname(__file__), '..',
                           'projects', '260324_신사업_운영_sop')
    out_path = os.path.join(out_dir, '신규_지역_진출_기획_sop.pptx')
    prs.save(out_path)
    print(f"PPTX saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
