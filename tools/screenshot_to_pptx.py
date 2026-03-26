"""
HTML 슬라이드를 Playwright로 캡처 → PPTX 이미지 슬라이드로 변환
"""
import os, asyncio, tempfile, shutil
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

SLIDE_W_PX = 1280
SLIDE_H_PX = 720
SLIDE_W_EMU = Emu(12192000)
SLIDE_H_EMU = Emu(6858000)

HTML_PATH = os.path.join(
    os.path.dirname(__file__), '..', 'projects',
    '260324_신사업_운영_sop', '신규_지역_진출_기획_sop_전체_v4_scroll.html'
)
OUT_PATH = os.path.join(
    os.path.dirname(__file__), '..', 'projects',
    '260324_신사업_운영_sop', '신규_지역_진출_기획_sop_이미지.pptx'
)


async def capture_slides(html_path, tmp_dir):
    file_url = 'file:///' + os.path.abspath(html_path).replace('\\', '/')
    images = []

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={'width': SLIDE_W_PX, 'height': SLIDE_H_PX})
        await page.goto(file_url, wait_until='networkidle')
        await page.wait_for_timeout(2000)

        slides = await page.query_selector_all('.slide-wrapper')
        print(f"Found {len(slides)} slides")

        for i, slide in enumerate(slides):
            img_path = os.path.join(tmp_dir, f'slide_{i+1:02d}.png')
            await slide.screenshot(path=img_path)
            images.append(img_path)
            print(f"  Captured slide {i+1}")

        await browser.close()

    return images


def build_pptx(images, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU

    for img_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_W_EMU, SLIDE_H_EMU)

    prs.save(out_path)
    print(f"\nPPTX saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


async def main():
    html = os.path.abspath(HTML_PATH)
    out = os.path.abspath(OUT_PATH)
    print(f"HTML: {html}")
    print(f"Output: {out}\n")

    tmp_dir = tempfile.mkdtemp(prefix='sop_slides_')
    try:
        images = await capture_slides(html, tmp_dir)
        build_pptx(images, out)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


if __name__ == '__main__':
    asyncio.run(main())
