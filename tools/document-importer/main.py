#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Document Importer
=================
PDF, PPTX, PPT 파일을 텍스트 + 이미지로 추출하는 도구

사용법:
    python main.py <파일경로>
    
예시:
    python main.py inbox/report.pdf
    python main.py inbox/presentation.pptx
"""

import sys
import os
import shutil
import io
from pathlib import Path
from datetime import datetime
from typing import Optional, Tuple

# Windows 콘솔 UTF-8 인코딩 설정
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# 프로젝트 루트 경로 (tools/document-importer 기준으로 2단계 상위)
SCRIPT_DIR = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent


def get_file_type(file_path: Path) -> Optional[str]:
    """파일 확장자로 타입 감지"""
    ext = file_path.suffix.lower()
    if ext == '.pdf':
        return 'pdf'
    elif ext in ['.pptx', '.ppt']:
        return 'pptx'
    return None


def create_temp_folder(file_path: Path) -> Path:
    """임시 작업 폴더 생성"""
    file_name = file_path.stem
    temp_dir = file_path.parent / f"_temp_{file_name}"
    
    # 기존 폴더가 있으면 삭제
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    
    # 폴더 구조 생성
    (temp_dir / "pages_text").mkdir(parents=True)
    (temp_dir / "pages_images").mkdir(parents=True)
    
    return temp_dir


def extract_pdf_text(pdf_path: Path, output_dir: Path) -> int:
    """
    PDF 텍스트 추출 (페이지별)
    
    Returns:
        추출된 페이지 수
    """
    try:
        import pdfplumber
    except ImportError:
        print("[오류] pdfplumber가 설치되지 않았습니다.")
        print("       pip install pdfplumber 실행 후 다시 시도하세요.")
        return 0
    
    text_dir = output_dir / "pages_text"
    all_text = []
    page_count = 0
    
    with pdfplumber.open(pdf_path) as pdf:
        page_count = len(pdf.pages)
        
        for i, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text() or ""
            
            # 개별 페이지 저장
            page_file = text_dir / f"page_{i:02d}.md"
            with open(page_file, 'w', encoding='utf-8') as f:
                f.write(f"# 페이지 {i}\n\n")
                f.write(page_text)
            
            all_text.append(f"## 페이지 {i}\n\n{page_text}")
            print(f"  [텍스트] 페이지 {i}/{page_count} 추출 완료")
    
    # 전체 텍스트 병합본 저장
    all_text_file = output_dir / "all_text.md"
    with open(all_text_file, 'w', encoding='utf-8') as f:
        f.write(f"# {pdf_path.stem} - 전체 텍스트\n\n")
        f.write(f"> 추출일: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n")
        f.write(f"> 페이지 수: {page_count}\n\n")
        f.write("---\n\n")
        f.write("\n\n---\n\n".join(all_text))
    
    return page_count


def extract_pptx_text(pptx_path: Path, output_dir: Path) -> int:
    """
    PPTX 텍스트 추출 (슬라이드별)
    
    Returns:
        추출된 슬라이드 수
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("[오류] python-pptx가 설치되지 않았습니다.")
        print("       pip install python-pptx 실행 후 다시 시도하세요.")
        return 0
    
    text_dir = output_dir / "pages_text"
    all_text = []
    
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # 슬라이드의 모든 shape에서 텍스트 추출
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        page_content = "\n\n".join(slide_text)
        
        # 개별 슬라이드 저장
        page_file = text_dir / f"page_{i:02d}.md"
        with open(page_file, 'w', encoding='utf-8') as f:
            f.write(f"# 슬라이드 {i}\n\n")
            f.write(page_content)
        
        all_text.append(f"## 슬라이드 {i}\n\n{page_content}")
        print(f"  [텍스트] 슬라이드 {i}/{slide_count} 추출 완료")
    
    # 전체 텍스트 병합본 저장
    all_text_file = output_dir / "all_text.md"
    with open(all_text_file, 'w', encoding='utf-8') as f:
        f.write(f"# {pptx_path.stem} - 전체 텍스트\n\n")
        f.write(f"> 추출일: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n")
        f.write(f"> 슬라이드 수: {slide_count}\n\n")
        f.write("---\n\n")
        f.write("\n\n---\n\n".join(all_text))
    
    return slide_count


def convert_pdf_to_images(pdf_path: Path, output_dir: Path, dpi: int = 150) -> bool:
    """
    PDF → 이미지 변환 (페이지별)
    
    Args:
        dpi: 이미지 해상도 (기본 150, 높을수록 선명하지만 용량 증가)
    
    Returns:
        성공 여부
    """
    try:
        from pdf2image import convert_from_path
    except ImportError:
        print("[오류] pdf2image가 설치되지 않았습니다.")
        print("       pip install pdf2image 실행 후 다시 시도하세요.")
        return False
    
    # Poppler 경로 확인 (Windows)
    poppler_path = None
    
    # 일반적인 Poppler 설치 경로들
    possible_paths = [
        Path(os.environ.get('POPPLER_PATH', '')) if os.environ.get('POPPLER_PATH') else None,
        Path("C:/Program Files/poppler/Library/bin"),
        Path("C:/Program Files/poppler-24.02.0/Library/bin"),
        Path("C:/poppler/Library/bin"),
        Path.home() / "poppler" / "Library" / "bin",
    ]
    
    for path in possible_paths:
        if path and path.exists():
            poppler_path = str(path)
            break
    
    image_dir = output_dir / "pages_images"
    
    try:
        if poppler_path:
            images = convert_from_path(pdf_path, dpi=dpi, poppler_path=poppler_path)
        else:
            # PATH에 있다고 가정
            images = convert_from_path(pdf_path, dpi=dpi)
        
        for i, image in enumerate(images, 1):
            image_file = image_dir / f"page_{i:02d}.png"
            image.save(image_file, 'PNG')
            print(f"  [이미지] 페이지 {i}/{len(images)} 변환 완료")
        
        return True
        
    except Exception as e:
        print(f"[오류] PDF 이미지 변환 실패: {e}")
        print("\n[해결 방법]")
        print("  1. Poppler가 설치되어 있는지 확인하세요")
        print("  2. setup.ps1을 실행하여 Poppler를 설치하세요")
        print("  3. 또는 환경변수 POPPLER_PATH를 설정하세요")
        return False


def convert_pptx_to_images(pptx_path: Path, output_dir: Path) -> bool:
    """
    PPTX → 이미지 변환 (PowerPoint COM 활용)
    
    Windows + PowerPoint 설치 필요
    
    Returns:
        성공 여부
    """
    # 절대 경로로 변환
    pptx_path = pptx_path.resolve()
    image_dir = (output_dir / "pages_images").resolve()
    
    try:
        import comtypes.client
    except ImportError:
        print("[오류] comtypes가 설치되지 않았습니다.")
        print("       pip install comtypes 실행 후 다시 시도하세요.")
        return False
    
    try:
        # PowerPoint 애플리케이션 실행
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # 창 표시 (안정성 위해)
        
        # 프레젠테이션 열기
        presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        
        slide_count = len(presentation.Slides)
        
        for i, slide in enumerate(presentation.Slides, 1):
            image_file = image_dir / f"page_{i:02d}.png"
            slide.Export(str(image_file), "PNG", 1920, 1080)
            print(f"  [이미지] 슬라이드 {i}/{slide_count} 변환 완료")
        
        presentation.Close()
        powerpoint.Quit()
        
        return True
        
    except Exception as e:
        print(f"[오류] PPTX 이미지 변환 실패: {e}")
        print("\n[해결 방법]")
        print("  1. Microsoft PowerPoint가 설치되어 있는지 확인하세요")
        print("  2. PowerPoint가 다른 프로세스에서 사용 중이 아닌지 확인하세요")
        return False


def print_completion_message(file_path: Path, temp_dir: Path, page_count: int, 
                              text_success: bool, image_success: bool):
    """완료 메시지 출력"""
    
    print("\n" + "=" * 60)
    print("추출 완료!")
    print("=" * 60)
    
    print(f"\n원본 파일: {file_path}")
    print(f"페이지 수: {page_count}")
    print(f"작업 폴더: {temp_dir}")
    
    print("\n[추출 결과]")
    print(f"  텍스트 추출: {'성공' if text_success else '실패'}")
    print(f"  이미지 변환: {'성공' if image_success else '실패 (텍스트만 사용)'}")
    
    if image_success:
        print("\n[다음 단계]")
        print("  Cursor AI에게 다음과 같이 요청하세요:")
        print("")
        print('  "@' + str(temp_dir.relative_to(PROJECT_ROOT)) + ' 분석해줘"')
        print("")
        print("  AI가 이미지를 분석하고 종합 마크다운을 생성합니다.")
        print("  리뷰 후 '확정' 또는 '정리해줘'라고 하면 저장됩니다.")
    else:
        print("\n[다음 단계]")
        print("  이미지 변환이 실패하여 텍스트만 추출되었습니다.")
        print("  Cursor AI에게 다음과 같이 요청하세요:")
        print("")
        print(f'  "{temp_dir / "all_text.md"} 분석하고 정리해줘"')
    
    print("\n" + "=" * 60)


def process_document(file_path_str: str) -> Tuple[bool, Optional[Path]]:
    """
    메인 처리 함수
    
    Args:
        file_path_str: 처리할 파일 경로
        
    Returns:
        (성공 여부, 임시 폴더 경로)
    """
    # 경로 처리
    file_path = Path(file_path_str)
    
    # 상대 경로면 프로젝트 루트 기준으로 변환
    if not file_path.is_absolute():
        file_path = PROJECT_ROOT / file_path
    
    # 파일 존재 확인
    if not file_path.exists():
        print(f"[오류] 파일을 찾을 수 없습니다: {file_path}")
        return False, None
    
    # 파일 타입 확인
    file_type = get_file_type(file_path)
    if not file_type:
        print(f"[오류] 지원하지 않는 파일 형식입니다: {file_path.suffix}")
        print("       지원 형식: PDF, PPTX, PPT")
        return False, None
    
    print(f"\n[Document Importer]")
    print(f"파일: {file_path.name}")
    print(f"타입: {file_type.upper()}")
    print("-" * 40)
    
    # 임시 폴더 생성
    temp_dir = create_temp_folder(file_path)
    print(f"작업 폴더 생성: {temp_dir.name}")
    
    # 텍스트 추출
    print("\n[1/2] 텍스트 추출 중...")
    if file_type == 'pdf':
        page_count = extract_pdf_text(file_path, temp_dir)
    else:
        page_count = extract_pptx_text(file_path, temp_dir)
    
    text_success = page_count > 0
    
    # 이미지 변환
    print("\n[2/2] 이미지 변환 중...")
    if file_type == 'pdf':
        image_success = convert_pdf_to_images(file_path, temp_dir)
    else:
        image_success = convert_pptx_to_images(file_path, temp_dir)
    
    # 완료 메시지
    print_completion_message(file_path, temp_dir, page_count, text_success, image_success)
    
    return text_success or image_success, temp_dir


def find_files_in_inbox() -> list:
    """inbox/ 폴더에서 처리 가능한 파일 목록 반환"""
    inbox_dir = PROJECT_ROOT / "inbox"
    if not inbox_dir.exists():
        return []
    
    files = []
    for ext in ['.pdf', '.pptx', '.ppt']:
        files.extend(inbox_dir.glob(f'*{ext}'))
    
    return sorted(files)


def main():
    """CLI 엔트리포인트"""
    
    # 인자 없이 실행하면 inbox/ 폴더 스캔
    if len(sys.argv) < 2:
        files = find_files_in_inbox()
        
        if not files:
            print(__doc__)
            print("\n[사용 예시]")
            print("  python main.py inbox/report.pdf")
            print("  python main.py inbox/presentation.pptx")
            print("  python main.py resources/document.pdf")
            print("\n[자동 검색]")
            print("  인자 없이 실행하면 inbox/ 폴더를 스캔합니다.")
            print("  현재 inbox/에 처리 가능한 파일이 없습니다.")
            sys.exit(1)
        
        print("\n[inbox/ 폴더 스캔 결과]")
        print("-" * 40)
        for i, f in enumerate(files, 1):
            print(f"  {i}. {f.name}")
        print("-" * 40)
        
        if len(files) == 1:
            print(f"\n파일 1개 발견. 자동 처리합니다: {files[0].name}")
            file_path = str(files[0])
        else:
            try:
                choice = input(f"\n처리할 파일 번호를 입력하세요 (1-{len(files)}): ").strip()
                idx = int(choice) - 1
                if 0 <= idx < len(files):
                    file_path = str(files[idx])
                else:
                    print("[오류] 잘못된 번호입니다.")
                    sys.exit(1)
            except (ValueError, EOFError):
                print("[오류] 숫자를 입력해주세요.")
                sys.exit(1)
    else:
        file_path = sys.argv[1]
    
    success, temp_dir = process_document(file_path)
    
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
